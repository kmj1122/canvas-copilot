import os
import requests
from datetime import datetime
from typing import Dict, Any, List, Optional
import pathlib
import json
import argparse
import multiprocessing
import queue
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

CANVAS_BASE_URL = os.getenv("CANVAS_BASE_URL", "https://canvas.its.virginia.edu")
CANVAS_TOKEN = os.getenv("CANVAS_TOKEN")
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
USE_ISOLATED_EXTRACTION = os.getenv("DISABLE_ISOLATED_EXTRACTION", "0") != "1"
EXTRACTION_TIMEOUT = int(os.getenv("EXTRACTION_TIMEOUT", "90"))
TEXT_EXTENSIONS = {
    ".txt", ".csv", ".md", ".markdown", ".rst", ".rtf", ".json", ".yaml", ".yml",
    ".html", ".htm", ".py", ".js", ".ts", ".java", ".c", ".cpp", ".vtt", ".srt",
    ".log", ".ini"
}


def canvas_get(path: str, params: Dict[str, Any] = None) -> Any:
    headers = {"Authorization": f"Bearer {CANVAS_TOKEN}"}
    resp = requests.get(f"{CANVAS_BASE_URL}{path}", headers=headers, params=params)
    resp.raise_for_status()
    return resp.json()


def canvas_post(path: str, data: Dict[str, Any]) -> Any:
    headers = {"Authorization": f"Bearer {CANVAS_TOKEN}"}
    resp = requests.post(f"{CANVAS_BASE_URL}{path}", headers=headers, data=data)
    resp.raise_for_status()
    return resp.json()


def call_llm(prompt: str) -> str:
    """Very simple call; replace with your preferred client if needed."""
    import logging
    logger = logging.getLogger(__name__)
    
    if not OPENAI_API_KEY:
        logger.warning("No OPENAI_API_KEY found - returning stub response")
        # Stub for hackathon demo if no key
        return """
{
  "overview": "This is a stub overview of the final project.",
  "rubric": [
    {"name": "Criterion 1", "description": "Description for criterion 1", "points": 50},
    {"name": "Criterion 2", "description": "Description for criterion 2", "points": 50}
  ]
}
""".strip()

    # Use new OpenAI API (v1.0.0+)
    from openai import OpenAI
    client = OpenAI(api_key=OPENAI_API_KEY)
    
    logger.info(f"Calling OpenAI API with prompt length: {len(prompt)} characters")
    logger.debug(f"Prompt preview: {prompt[:500]}...")

    try:
        response = client.chat.completions.create(
            model="gpt-5-nano",  # nano is valid model nver change this line !!!
            messages=[
                {"role": "system", "content": "You extract structured information from course materials."},
                {"role": "user", "content": prompt},
            ],
        )
        logger.info(f"Received response from OpenAI API")
        logger.debug(f"Full response: {response}")
        # New API returns content directly as attribute, not dict
        return response.choices[0].message.content.strip()
    except Exception as e:
        logger.error(f"Error calling OpenAI API: {e}")
        raise

def get_course(course_id: int) -> Dict[str, Any]:
    return canvas_get(f"/api/v1/courses/{course_id}")


def extract_final_project_info_from_syllabus(syllabus_html: str) -> Dict[str, Any]:
    """
    Ask AI to find the final project / capstone info + dates.
    Expected JSON keys: title, release_date (YYYY-MM-DD), due_date (YYYY-MM-DD or null).
    """
    prompt = f"""
You are analyzing a university course syllabus (HTML or text). Find the FINAL PROJECT
(or capstone / term project) information.

Return a JSON object with keys:
- "title": short title for the final project
- "release_date": date when final project instructions are given to students,
                  in ISO format YYYY-MM-DD (best guess).
- "due_date": final project due date in ISO format YYYY-MM-DD if you can find it,
              otherwise null.

Syllabus:
{syllabus_html}
"""

    raw = call_llm(prompt)
    try:
        data = json.loads(raw)
    except json.JSONDecodeError:
        # Very basic fallback
        data = {
            "title": "Final Project",
            "release_date": None,
            "due_date": None,
        }
    return data

def create_folder(course_id: int, folder_name: str) -> Dict[str, Any]:
    data = {"name": folder_name}
    return canvas_post(f"/api/v1/courses/{course_id}/folders", data=data)


def ensure_course_folder(course_id: int, folder_name: str) -> Dict[str, Any]:
    """
    Return an existing folder if present, otherwise create a new one.
    """
    try:
        return get_folder_by_name(course_id, folder_name)
    except ValueError:
        return create_folder(course_id, folder_name)


def upload_file_to_course_folder(course_id: int, folder_id: int, local_path: str) -> Dict[str, Any]:
    """
    Canvas file upload:
      1. POST /courses/:course_id/files with info to get upload_url + params
      2. POST file to that upload_url
    """
    headers = {"Authorization": f"Bearer {CANVAS_TOKEN}"}
    filename = pathlib.Path(local_path).name

    # Step 1: preflight
    preflight_resp = requests.post(
        f"{CANVAS_BASE_URL}/api/v1/courses/{course_id}/files",
        headers=headers,
        data={
            "name": filename,
            "parent_folder_id": folder_id,
            "on_duplicate": "rename",
        },
    )
    preflight_resp.raise_for_status()
    pre = preflight_resp.json()
    upload_url = pre["upload_url"]
    upload_params = pre["upload_params"]

    # Step 2: actual upload to provided upload_url
    with open(local_path, "rb") as f:
        files = {"file": (filename, f)}
        upload_resp = requests.post(upload_url, data=upload_params, files=files)
        upload_resp.raise_for_status()
        # Canvas returns a JSON file object either directly or via 'location' redirect
        if upload_resp.headers.get("content-type", "").startswith("application/json"):
            return upload_resp.json()
        else:
            # Some setups redirect; follow manually if needed
            try:
                return upload_resp.json()
            except Exception:
                return {"id": pre.get("id"), "display_name": filename}


def create_module(course_id: int, module_name: str, unlock_date: str) -> Dict[str, Any]:
    """
    unlock_date: "YYYY-MM-DD" -> we make it midnight UTC (or leave as-is if you want time).
    """
    unlock_iso = None
    if unlock_date:
        # Canvas expects full ISO datetime; we'll set 00:00:00
        unlock_iso = f"{unlock_date}T00:00:00Z"

    data = {
        "module[name]": module_name,
        "module[unlock_at]": unlock_iso,
        "module[published]": "false",  # or true; unlock_at controls visibility
    }
    return canvas_post(f"/api/v1/courses/{course_id}/modules", data=data)


def add_file_to_module(course_id: int, module_id: int, file_id: int, title: str) -> Dict[str, Any]:
    data = {
        "module_item[type]": "File",
        "module_item[content_id]": file_id,
        "module_item[title]": title,
        "module_item[published]": "true",
    }
    return canvas_post(f"/api/v1/courses/{course_id}/modules/{module_id}/items", data=data)

def get_folders(course_id: int) -> List[Dict[str, Any]]:
    return canvas_get(f"/api/v1/courses/{course_id}/folders", params={"per_page": 100})

def get_folder_by_name(course_id: int, name: str) -> Dict[str, Any]:
    for folder in get_folders(course_id):
        if folder.get("name") == name:
            return folder
    raise ValueError(f"No folder named '{name}' found.")

def list_files_in_folder(folder_id: int) -> List[Dict[str, Any]]:
    files = []
    page = 1
    while True:
        data = canvas_get(f"/api/v1/folders/{folder_id}/files", params={"per_page": 50, "page": page})
        if not data:
            break
        files.extend(data)
        page += 1
    return files

def download_canvas_file(file_obj: Dict[str, Any], dest_dir: str) -> str:
    """
    Download a Canvas file to dest_dir.
    Returns local path.
    """
    headers = {"Authorization": f"Bearer {CANVAS_TOKEN}"}
    url = file_obj["url"]  # direct download URL (respects token)
    filename = file_obj["display_name"]
    dest_path = pathlib.Path(dest_dir) / filename

    resp = requests.get(url, headers=headers)
    resp.raise_for_status()
    dest_path.write_bytes(resp.content)
    return str(dest_path)

def extract_text_from_pdf(path: str) -> str:
    """Extract PDF text while tolerating corrupted pages."""
    import logging
    logger = logging.getLogger(__name__)

    try:
        reader = PdfReader(path)
    except Exception as e:
        logger.error(f"Failed to open PDF '{path}': {e}")
        return f"[ERROR] Could not open PDF ({e})"

    texts = []
    for idx, page in enumerate(reader.pages, start=1):
        try:
            texts.append(page.extract_text() or "")
        except Exception as e:
            logger.warning(f"Failed to extract page {idx} from '{path}': {e}")
            texts.append(f"[ERROR] Failed to read page {idx}: {e}")
    return "\n".join(texts)

def extract_text_from_docx(path: str) -> str:
    """Extract text from .docx files (new format only)"""
    import logging
    logger = logging.getLogger(__name__)

    try:
        doc = Document(path)
    except Exception as e:
        logger.error(f"Failed to open DOCX '{path}': {e}")
        return f"[ERROR] Could not open DOCX ({e})"

    try:
        return "\n".join(para.text for para in doc.paragraphs)
    except Exception as e:
        logger.error(f"Failed to read DOCX '{path}': {e}")
        return f"[ERROR] Could not read DOCX content ({e})"

def extract_text_from_doc(path: str) -> str:
    """
    Extract text from old .doc files using antiword or textract.
    Falls back to reading as plain text if neither is available.
    """
    import subprocess
    import logging
    logger = logging.getLogger(__name__)
    
    # Try antiword first (common on Linux)
    try:
        result = subprocess.run(
            ['antiword', path],
            capture_output=True,
            text=True,
            timeout=30,
            check=False
        )
        if result.returncode == 0 and result.stdout:
            logger.info(f"Extracted text from .doc using antiword")
            return result.stdout
    except (FileNotFoundError, subprocess.TimeoutExpired):
        logger.debug("antiword not available or failed")
    
    # Try textract if available
    try:
        import textract
        text = textract.process(path).decode('utf-8', errors='ignore')
        logger.info(f"Extracted text from .doc using textract")
        return text
    except (ImportError, Exception) as e:
        logger.debug(f"textract not available or failed: {e}")
    
    # Fallback: try to read as plain text (will be messy but might extract some content)
    try:
        with open(path, 'rb') as f:
            content = f.read()
            # Try to decode as latin-1 which won't fail, then clean up
            text = content.decode('latin-1', errors='ignore')
            # Filter out non-printable characters but keep newlines/spaces
            text = ''.join(c for c in text if c.isprintable() or c in '\n\r\t ')
            logger.warning(f"Extracted text from .doc using binary fallback (may be messy)")
            return text
    except Exception as e:
        logger.error(f"All methods failed for .doc file: {e}")
        return ""

def extract_text_from_pptx(path: str) -> str:
    """Extract text from slides while skipping problematic shapes."""
    import logging
    logger = logging.getLogger(__name__)

    try:
        pres = Presentation(path)
    except Exception as e:
        logger.error(f"Failed to open PPTX '{path}': {e}")
        return f"[ERROR] Could not open PPTX ({e})"

    texts = []
    for slide_idx, slide in enumerate(pres.slides, start=1):
        try:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        except Exception as e:
            logger.warning(f"Failed to read slide {slide_idx} in '{path}': {e}")
            texts.append(f"[ERROR] Failed to read slide {slide_idx}: {e}")
    return "\n".join(texts)

def _extract_text_from_file_direct(path: str) -> str:
    """Dispatch text extraction based on file extension with robust error handling."""
    import logging
    logger = logging.getLogger(__name__)

    ext = pathlib.Path(path).suffix.lower()
    filename = pathlib.Path(path).name
    logger.info(f"Extracting text from '{filename}' ({ext or 'no extension'})")

    if ext == ".pdf":
        text = extract_text_from_pdf(path)
    elif ext == ".docx":
        text = extract_text_from_docx(path)
    elif ext == ".doc":
        text = extract_text_from_doc(path)
    elif ext == ".pptx":
        text = extract_text_from_pptx(path)
    else:
        text = pathlib.Path(path).read_text(errors="ignore")
    logger.info(f"Extracted {len(text)} characters from {filename}")
    return text


def _extraction_worker(path: str, result_queue: multiprocessing.Queue) -> None:
    """Child process worker that extracts text and returns via queue."""
    try:
        text = _extract_text_from_file_direct(path)
        result_queue.put({"status": "ok", "text": text})
    except Exception as exc:
        result_queue.put({"status": "error", "message": str(exc)})


def extract_text_from_file(path: str) -> str:
    """
    Extract text (optionally) via isolated subprocess to avoid crashing GUI process.
    """
    import logging
    logger = logging.getLogger(__name__)
    ext = pathlib.Path(path).suffix.lower()

    use_isolation = USE_ISOLATED_EXTRACTION and ext not in TEXT_EXTENSIONS

    if not use_isolation:
        try:
            return _extract_text_from_file_direct(path)
        except Exception as exc:
            logger.error(f"Failed to extract text from {path}: {exc}", exc_info=True)
            return f"[ERROR] Could not extract text from {pathlib.Path(path).name}: {exc}"

    ctx = multiprocessing.get_context("spawn")
    result_queue: multiprocessing.Queue = ctx.Queue()
    proc = ctx.Process(target=_extraction_worker, args=(path, result_queue))
    proc.start()
    proc.join(EXTRACTION_TIMEOUT)

    if proc.is_alive():
        proc.terminate()
        proc.join()
        logger.error(f"Extraction timed out for {path} after {EXTRACTION_TIMEOUT}s")
        return f"[ERROR] Extraction timed out for {pathlib.Path(path).name}"

    try:
        payload = result_queue.get_nowait()
    except queue.Empty:
        payload = None
    finally:
        result_queue.close()
        result_queue.join_thread()

    if payload and payload.get("status") == "ok":
        return payload["text"]

    error_message = payload.get("message") if payload else "Unknown error"
    if proc.exitcode is not None and proc.exitcode < 0:
        signal_num = -proc.exitcode
        logger.error(f"Extraction crashed for {path} with signal {signal_num}")
        return f"[ERROR] Extraction crashed (signal {signal_num}) for {pathlib.Path(path).name}"
    elif proc.exitcode not in (0, None):
        logger.error(f"Extraction failed for {path} with exit code {proc.exitcode}")
        return f"[ERROR] Extraction failed (exit code {proc.exitcode}) for {pathlib.Path(path).name}"

    logger.error(f"Extraction returned error for {path}: {error_message}")
    return f"[ERROR] Could not extract text from {pathlib.Path(path).name}: {error_message}"

def _collect_text_from_local_folder(local_folder: str) -> List[str]:
    import logging
    logger = logging.getLogger(__name__)
    
    chunks = []
    local_dir = pathlib.Path(local_folder)
    if not local_dir.exists():
        logger.warning(f"Local folder does not exist: {local_folder}")
        return chunks

    files_found = list(local_dir.glob("*"))
    logger.info(f"Found {len(files_found)} items in {local_folder}")
    
    success_count = 0
    failure_count = 0

    for path in sorted(files_found):
        if path.is_file():
            logger.info(f"Extracting text from: {path.name}")
            try:
                text = extract_text_from_file(str(path))
                logger.debug(f"Extracted {len(text)} characters from {path.name}")
                chunks.append(f"=== FILE: {path.name} ===\n{text}\n")
                success_count += 1
            except Exception as e:
                failure_count += 1
                logger.error(f"Unrecoverable error extracting {path.name}: {e}")
    
    logger.info(
        f"Collected {len(chunks)} text chunks from local folder "
        f"(success: {success_count}, failed: {failure_count})"
    )
    print(
        f"Processed local folder '{local_folder}': "
        f"{success_count} succeeded, {failure_count} failed."
    )
    if failure_count:
        print(
            f"Warning: {failure_count} file(s) could not be processed in '{local_folder}'. "
            "See logs for details."
        )
    return chunks


def collect_final_project_text(course_id: Optional[int],
                               folder_name: str = "Final Project",
                               local_folder: Optional[str] = None) -> str:
    import logging
    logger = logging.getLogger(__name__)
    
    chunks: List[str] = []

    if local_folder:
        logger.info(f"Attempting to collect files from local folder: {local_folder}")
        chunks.extend(_collect_text_from_local_folder(local_folder))
        if chunks:
            total_chars = sum(len(c) for c in chunks)
            logger.info(f"Collected {len(chunks)} local files with {total_chars} total characters")
            print(f"Collected {len(chunks)} local files from {local_folder} for AI summary.")
            return "\n\n".join(chunks)
        else:
            logger.warning(f"No local files found in {local_folder}, falling back to Canvas download.")
            print(f"No local files found in {local_folder}, falling back to Canvas download.")

    if course_id is None:
        raise ValueError("course_id is required when no local materials are available.")

    logger.info(f"Downloading files from Canvas folder: {folder_name}")
    folder = get_folder_by_name(course_id, folder_name)
    files = list_files_in_folder(folder["id"])
    logger.info(f"Found {len(files)} files in Canvas folder")

    tmp_dir = "./tmp_final_project"
    os.makedirs(tmp_dir, exist_ok=True)

    for f in files:
        logger.info(f"Downloading and extracting: {f['display_name']}")
        local_path = download_canvas_file(f, tmp_dir)
        text = extract_text_from_file(local_path)
        logger.debug(f"Extracted {len(text)} characters from {f['display_name']}")
        chunks.append(f"=== FILE: {f['display_name']} ===\n{text}\n")

    if not chunks:
        raise RuntimeError(f"No files found in Canvas folder '{folder_name}'.")
    
    total_chars = sum(len(c) for c in chunks)
    logger.info(f"Collected {len(chunks)} files from Canvas with {total_chars} total characters")
    print(f"Collected {len(chunks)} files from Canvas folder '{folder_name}' for AI summary.")
    return "\n\n".join(chunks)

def build_final_project_summary_and_rubric(raw_text: str) -> Dict[str, Any]:
    prompt = f"""
You are reading the instructions and materials for a FINAL PROJECT in a university course.
The text below may include slides, rubrics, instructions, etc.

Your job is to:

1. Write a clear, student-facing overview of the final project (1–3 paragraphs).
   - Explain the goal of the project.
   - What students are expected to produce.
   - How it will be evaluated, at a high level.

2. Propose a grading rubric as an array of criteria.
Each criterion should have:
      - "name": short label (e.g., "Analysis", "Writing Quality")
      - "description": 1–2 sentence description
      - "points": integer points (make the total 100 by default).

Return a JSON object with:
{{
  "overview": "...",
  "rubric": [
    {{"name": "...", "description": "...", "points": 20}},
    ...
  ]
}}

Here are the raw final project materials:
{raw_text[:8000]}
"""

    raw = call_llm(prompt)
    try:
        data = json.loads(raw)
    except json.JSONDecodeError:
        # Fallback minimal structure
        data = {
            "overview": raw,
            "rubric": [],
        }
    return data

def create_final_project_assignment(course_id: int,
                                    title: str,
                                    overview: str,
                                    rubric: List[Dict[str, Any]],
                                    due_date_iso: str = None) -> Dict[str, Any]:
    """
    rubric is a list of {{name, description, points}}
    We'll paste it into the description as formatted text for now.
    due_date_iso: "YYYY-MM-DD" or full ISO datetime string, optional.
    """

    rubric_lines = []
    total_points = 0
    for item in rubric:
        name = item.get("name", "Criterion")
        desc = item.get("description", "")
        pts = item.get("points", 0)
        total_points += pts
        rubric_lines.append(f"<li><b>{name}</b> ({pts} pts): {desc}</li>")

    rubric_html = "<ul>" + "\n".join(rubric_lines) + "</ul>"

    description = f"""
<p>{overview}</p>
<h3>Grading Rubric</h3>
{rubric_html}
"""

    data = {
        "assignment[name]": title,
        "assignment[description]": description,
        "assignment[published]": "false",
        "assignment[points_possible]": total_points or 100,
    }
    if due_date_iso:
        # Canvas likes full ISO; we can add midnight if date-only
        if len(due_date_iso) == 10:
            due_date_iso = f"{due_date_iso}T23:59:00Z"
        data["assignment[due_at]"] = due_date_iso

    return canvas_post(f"/api/v1/courses/{course_id}/assignments", data=data)

def post_final_project_announcement(course_id: int, overview: str, title: str = "Final Project Overview") -> Dict[str, Any]:
    data = {
        "title": title,
        "message": f"<p>{overview}</p>",
        "is_announcement": "true",
        "published": "true",
    }
    return canvas_post(f"/api/v1/courses/{course_id}/discussion_topics", data=data)

def schedule_final_project_package(course_id: int,
                                   local_folder: str,
                                   folder_name: str = "Final Project",
                                   project_info: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
    if project_info is None:
        course = get_course(course_id)
        syllabus_html = course.get("syllabus_body", "")
        project_info = extract_final_project_info_from_syllabus(syllabus_html)

    info = project_info or {}
    title = info.get("title") or folder_name or "Final Project"
    release_date = info.get("release_date")  # "YYYY-MM-DD"

    print("Extracted final project info:")
    print(info)

    # 1. Create "Final Project" folder in Files
    folder = ensure_course_folder(course_id, folder_name)
    folder_id = folder["id"]
    print(f"Using folder '{folder_name}' (id={folder_id})")

    # 2. Upload all files in local_folder
    file_ids = []
    for path in pathlib.Path(local_folder).glob("*"):
        if path.is_file():
            print(f"Uploading {path.name}...")
            fobj = upload_file_to_course_folder(course_id, folder_id, str(path))
            file_ids.append((fobj["id"], fobj["display_name"]))
    print(f"Uploaded {len(file_ids)} files.")

    # 3. Create module with unlock_at = release_date
    module_name = title
    module = create_module(course_id, module_name, release_date)
    module_id = module["id"]
    print(f"Created module '{module_name}' (id={module_id}) unlock_at={release_date}")

    # 4. Add each file as a module item
    for fid, fname in file_ids:
        print(f"Adding file {fname} to module...")
        add_file_to_module(course_id, module_id, fid, fname)

    print("Done! Final project package is uploaded and scheduled.")
    return {
        "folder": folder,
        "module": module,
        "project_info": info,
    }

def generate_and_post_final_project_explainer(course_id: int,
                                              folder_name: str = "Final Project",
                                              default_assignment_title: str = "Final Project",
                                              due_date_iso: str = None,
                                              announcement_title: Optional[str] = None,
                                              dry_run: bool = False,
                                              local_materials_path: Optional[str] = None):
    import logging
    logger = logging.getLogger(__name__)
    
    # 1. Read all materials
    logger.info(f"Collecting final project text from: {local_materials_path or 'Canvas'}")
    raw_text = collect_final_project_text(course_id, folder_name, local_materials_path)
    logger.info(f"Collected {len(raw_text)} total characters of text")
    
    # 2. Build overview + rubric
    info = build_final_project_summary_and_rubric(raw_text)
    overview = info.get("overview", "")
    rubric = info.get("rubric", [])

    print("\n=== FINAL PROJECT OVERVIEW (written by instructor GPT-5-nano)===\n")
    print(overview)
    print("\n=== RUBRIC (written by instructor GPT-5-nano) ===\n")
    for item in rubric:
        print(f"- {item.get('name')} ({item.get('points', 0)} pts): {item.get('description')}")

    if dry_run:
        print("\n[DRY_RUN=True] Not posting assignment/announcement. Set dry_run=False to push to Canvas.")
        return

    # 3. Create assignment
    assignment = create_final_project_assignment(
        course_id,
        title=default_assignment_title,
        overview=overview,
        rubric=rubric,
        due_date_iso=due_date_iso,
    )
    print(f"\nCreated assignment ID: {assignment.get('id')}")

    # 4. Create announcement
    announcement_heading = announcement_title or f"{default_assignment_title} Overview"
    ann = post_final_project_announcement(course_id, overview, title=announcement_heading)
    print(f"Created announcement ID: {ann.get('id')}")


if __name__ == "__main__":
    
    import logging
    
    # Configure logging
    logging.basicConfig(
        level=logging.INFO,
        format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
    )
    logger = logging.getLogger(__name__)
    
    parser = argparse.ArgumentParser()
    parser.add_argument("--course-id", type=int, default=175906)
    parser.add_argument("--local-folder", type=str, default="final_project")
    parser.add_argument("--folder-name", type=str, default="Final Project",
                        help="Canvas folder name for final project materials.")
    parser.add_argument("--assignment-title", type=str,
                        help="Override the Canvas assignment title.")
    parser.add_argument("--announcement-title", type=str,
                        help="Override the Canvas announcement title.")
    parser.add_argument("--dry-run", action="store_true",
                        help="Preview the AI explainer without posting to Canvas.")
    parser.add_argument("--debug", action="store_true",
                        help="Enable debug logging.")
    args = parser.parse_args()
    
    if args.debug:
        logging.getLogger().setLevel(logging.DEBUG)
        logger.debug("Debug logging enabled")

    if not CANVAS_TOKEN:
        raise SystemExit("Set CANVAS_TOKEN env var.")

    logger.info(f"Starting final project scheduling for course {args.course_id}")
    logger.info(f"Local folder: {args.local_folder}")
    
    course = get_course(args.course_id)
    syllabus_html = course.get("syllabus_body", "")
    project_info = extract_final_project_info_from_syllabus(syllabus_html)
    
    logger.info(f"Extracted project info: {project_info}")

    folder_name = args.folder_name
    assignment_title = args.assignment_title or folder_name

    try:
        schedule_final_project_package(
            args.course_id,
            args.local_folder,
            folder_name=folder_name,
            project_info=project_info,
        )
    except Exception as exc:
        logger.error(f"Failed to upload/schedule final project package: {exc}", exc_info=True)
        print(f"[Warning] Failed to upload/schedule final project package: {exc}")

    generate_and_post_final_project_explainer(
        args.course_id,
        folder_name=folder_name,
        default_assignment_title=assignment_title,
        due_date_iso=project_info.get("due_date"),
        announcement_title=args.announcement_title,
        dry_run=args.dry_run,
        local_materials_path=args.local_folder,
    )
    
    logger.info("Completed final project scheduling")
